"""
Generate the SentinalIQ project presentation (.pptx) following the
"Tech Cast Presentation slides.pdf" template structure:

  1. Title (Sponsored by ASIA Charitable Trust / School of Computing)
  2. Project Title
  3. Team Name
  4. Introduction
  5. Problem Statement
  6. Objectives
  7. Existing System
  8. Technologies Used
  9. Flow of Project (diagram)
  10. Main Logic (code)
  11. Future Scope
  12. Thank You

Design: dark cyberpunk SOC theme (deep navy + neon cyan / purple / pink),
matching the SentinalIQ product brand.
"""

import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

# ── Design tokens ────────────────────────────────────────────────────────────
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.7
CONTENT_W = SLIDE_W - 2 * MARGIN

BG      = RGBColor(0x0A, 0x0E, 0x1A)   # page background
BG2     = RGBColor(0x0F, 0x15, 0x24)   # alt panel
PANEL   = RGBColor(0x11, 0x1A, 0x2C)   # cards / panels
EDGE    = RGBColor(0x1E, 0x29, 0x3B)   # borders
CODE_BG = RGBColor(0x0B, 0x12, 0x22)   # code panels
CHIP_BG = RGBColor(0x10, 0x1A, 0x2E)

CYAN    = RGBColor(0x22, 0xD3, 0xEE)
PINK    = RGBColor(0xF4, 0x72, 0xB6)
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)
GREEN   = RGBColor(0x34, 0xD3, 0x99)
AMBER   = RGBColor(0xFB, 0xBF, 0x24)
RED     = RGBColor(0xF8, 0x71, 0x71)
TEXT    = RGBColor(0xE2, 0xE8, 0xF0)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)

FONT = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
BLANK = prs.slide_layouts[6]


# ── Low-level helpers ────────────────────────────────────────────────────────
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color=BG):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def _rpr(run):
    return run._r.get_or_add_rPr()


def set_spacing(run, spc):
    _rpr(run).set("spc", str(spc))


def tx(s, x, y, w, h, text, size, color, bold=False, font=FONT,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True,
       spacing=None, line_spacing=None, italic=False):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    if spacing:
        set_spacing(r, spacing)
    return box


def shape(s, kind, x, y, w, h, fill=PANEL, line=EDGE, line_w=1.0, radius=None):
    sp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    sp.shadow.inherit = False
    return sp


def line(s, x1, y1, x2, y2, color=CYAN, width=1.5, dash=False, head=False, tail=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                               Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if dash:
        try:
            c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            ln = c.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if head:
        ln = c.line._get_or_add_ln()
        el = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(el)
    if tail:
        ln = c.line._get_or_add_ln()
        el = ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(el)
    c.shadow.inherit = False
    return c


def content_header(s, kicker, title, accent=CYAN, title_size=30):
    tx(s, MARGIN, 0.5, 11.5, 0.3, kicker, 12.5, accent, bold=True, spacing=300)
    tx(s, MARGIN, 0.84, 11.9, 0.62, title, title_size, TEXT, bold=True)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, 1.56, 1.15, 0.055,
          fill=accent, line=None, radius=0.5)


def footer(s, num):
    tx(s, MARGIN, 7.13, 9.0, 0.25, "SentinalIQ · Enterprise SIEM / SOC Platform", 9, MUTED)
    tx(s, 12.15, 7.13, 0.5, 0.25, f"{num:02d}", 9, MUTED, align=PP_ALIGN.RIGHT)


def panel(s, x, y, w, h, fill=PANEL, line_c=EDGE, radius=0.05):
    return shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=line_c, radius=radius)


def bullets(s, x, y, w, items, size=15, gap=9, accent=CYAN, color=TEXT,
            glyph="\u25B8", line_spacing=1.1):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        text, bold = item if isinstance(item, tuple) else (item, False)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.line_spacing = line_spacing
        g = p.add_run()
        g.text = glyph + "  "
        g.font.size = Pt(size)
        g.font.bold = True
        g.font.color.rgb = accent
        g.font.name = FONT
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT


def chip(s, x, y, w, h, text, accent=CYAN):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=CHIP_BG, line=EDGE, radius=0.5)
    shape(s, MSO_SHAPE.OVAL, x + 0.18, y + h / 2 - 0.05, 0.10, 0.10, fill=accent, line=None)
    tx(s, x + 0.4, y, w - 0.5, h, text, 11.5, TEXT, anchor=MSO_ANCHOR.MIDDLE)


def stat_card(s, x, y, w, h, value, label, accent=CYAN):
    panel(s, x, y, w, h)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.02, y, 0.06, h, fill=accent, line=None, radius=0.5)
    tx(s, x + 0.25, y + 0.12, w - 0.4, 0.4, value, 22, accent, bold=True)
    tx(s, x + 0.25, y + 0.52, w - 0.4, 0.28, label, 10.5, MUTED)


# ── Syntax-highlighted code ──────────────────────────────────────────────────
CODE_COLORS = {
    "comment": RGBColor(0x64, 0x74, 0x8B),
    "string":  RGBColor(0xFB, 0xBF, 0x24),
    "number":  RGBColor(0xA7, 0x8B, 0xFA),
    "deco":    RGBColor(0xF4, 0x72, 0xB6),
    "keyword": RGBColor(0xC7, 0x92, 0xEA),
    "func":    RGBColor(0x67, 0xE8, 0xF9),
    "plain":   TEXT,
}

TOKEN_RE = re.compile(
    r"(?P<comment>\#.*)"
    r"|(?P<string>f?\"(?:\\.|[^\"\\])*\"|f?'(?:\\.|[^'\\])*')"
    r"|(?P<deco>@\w+)"
    r"|(?P<number>\b\d[\d_]*\.?\d*\b)"
    r"|(?P<ident>[A-Za-z_]\w*)"
    r"|(?P<rest>[^\s])"
    r"|(?P<space>\s+)",
    re.VERBOSE,
)

KEYWORDS = {
    "def", "return", "import", "from", "async", "await", "class", "for",
    "in", "if", "else", "with", "as", "try", "except", "while", "None",
    "True", "False", "not", "and", "or", "pass", "raise", "yield", "is",
    "lambda", "global", "del", "self",
}


def tokenize(line):
    out = []
    for m in TOKEN_RE.finditer(line):
        kind = m.lastgroup
        text = m.group()
        if kind == "ident":
            if text in KEYWORDS:
                kind = "keyword"
            elif line[m.end():].lstrip().startswith("("):
                kind = "func"
        if kind == "space":
            continue
        out.append((text, CODE_COLORS.get(kind, CODE_COLORS["plain"])))
    return out


def code_panel(s, x, y, w, h, filename, code, dot_color=CYAN, font_size=9.5):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=CODE_BG, line=EDGE, radius=0.035)
    shape(s, MSO_SHAPE.RECTANGLE, x + 0.02, y + 0.02, w - 0.04, 0.36,
          fill=RGBColor(0x16, 0x20, 0x36), line=None)
    shape(s, MSO_SHAPE.OVAL, x + 0.17, y + 0.15, 0.10, 0.10, fill=dot_color, line=None)
    tx(s, x + 0.38, y + 0.05, w - 0.5, 0.3, filename, 10, MUTED,
       font=MONO, anchor=MSO_ANCHOR.MIDDLE)
    box = s.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.5),
                               Inches(w - 0.36), Inches(h - 0.65))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln_ in enumerate(code.strip("\n").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        p.line_spacing = 1.0
        for text, color in tokenize(ln_):
            r = p.add_run()
            r.text = text
            r.font.name = MONO
            r.font.size = Pt(font_size)
            r.font.color.rgb = color


def flow_box(s, x, y, w, h, title, sub, accent=CYAN):
    panel(s, x, y, w, h)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.02, y + 0.06, 0.05, h - 0.12,
          fill=accent, line=None, radius=0.5)
    tx(s, x + 0.15, y + 0.18, w - 0.3, 0.35, title, 13, TEXT, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x + 0.15, y + 0.55, w - 0.3, 0.45, sub, 9.5, MUTED, align=PP_ALIGN.CENTER, line_spacing=1.05)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 · TITLE
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 0.09, SLIDE_H, fill=CYAN, line=None)
shape(s, MSO_SHAPE.RECTANGLE, 0.09, 0, 0.09, SLIDE_H, fill=PURPLE, line=None)
shape(s, MSO_SHAPE.RECTANGLE, 0.18, 0, 0.09, SLIDE_H, fill=PINK, line=None)

tx(s, 1.2, 0.62, 11.0, 0.3, "SPONSORED BY ASIA CHARITABLE TRUST", 13, MUTED, bold=True,
   align=PP_ALIGN.CENTER, spacing=400)
tx(s, 1.2, 0.98, 11.0, 0.3, "SCHOOL OF COMPUTING", 13, CYAN, bold=True,
   align=PP_ALIGN.CENTER, spacing=400)

for dx, c in ((6.42, CYAN), (6.62, PURPLE), (6.82, PINK)):
    shape(s, MSO_SHAPE.OVAL, dx, 1.52, 0.07, 0.07, fill=c, line=None)

tx(s, 1.0, 2.28, 11.33, 1.3, "SentinalIQ", 72, TEXT, bold=True, align=PP_ALIGN.CENTER)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.97, 3.72, 5.4, 0.58, fill=None, line=CYAN, line_w=1.2, radius=0.5)
tx(s, 3.97, 3.72, 5.4, 0.58, "ENTERPRISE SIEM / SOC PLATFORM", 14, CYAN, bold=True,
   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=200)
tx(s, 1.0, 4.62, 11.33, 0.4,
   "Real-time Monitoring  ·  Incident Response  ·  Threat Intelligence  ·  AI Copilot",
   13, MUTED, align=PP_ALIGN.CENTER)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 5.87, 5.3, 1.6, 0.05, fill=PURPLE, line=None, radius=0.5)
tx(s, 1.0, 6.62, 11.33, 0.3, "TECH CAST · PROJECT PRESENTATION", 11, MUTED,
   align=PP_ALIGN.CENTER, spacing=250)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 · PROJECT TITLE
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
content_header(s, "PROJECT TITLE", "SentinalIQ — Enterprise SIEM / SOC Platform")

tx(s, MARGIN, 1.9, CONTENT_W, 1.0,
   "A full-stack Security Information & Event Management (SIEM) / Security "
   "Operations Center (SOC) platform built with React 19, TypeScript, "
   "Python FastAPI and SQLite — bringing dashboards, log exploration, "
   "incident management, endpoint & WAF monitoring and threat "
   "intelligence into one real-time command center.",
   15, MUTED, wrap=True, line_spacing=1.2)

stats = [
    ("30+", "REST API endpoints", CYAN),
    ("6", "Security modules", PURPLE),
    ("3", "User roles · JWT", PINK),
    ("WS", "Real-time alerts", GREEN),
]
sw, gap = (CONTENT_W - 3 * 0.12) / 4, 0.12
for i, (v, l, a) in enumerate(stats):
    stat_card(s, MARGIN + i * (sw + gap), 3.15, sw, 0.92, v, l, a)

tx(s, MARGIN, 4.32, 8, 0.3, "KEY MODULES", 11, MUTED, bold=True, spacing=300)
modules = [
    ("Security Dashboard", CYAN), ("Log Explorer", PURPLE),
    ("Incident Management", PINK), ("EDR / XDR", GREEN),
    ("WAF Monitor", AMBER), ("Threat Intelligence", CYAN),
    ("Data Ingestion", PURPLE), ("AI Copilot", PINK),
]
cw, cg = (CONTENT_W - 3 * 0.14) / 4, 0.14
for i, (label, a) in enumerate(modules):
    chip(s, MARGIN + (i % 4) * (cw + cg), 4.72 + (i // 4) * 0.62, cw, 0.5, label, a)

footer(s, 2)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 · TEAM NAME
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
content_header(s, "TEAM NAME", "Meet the Team")

team = [
    ("N1", "Name 1", "Team Lead · Backend", PINK),
    ("N2", "Name 2", "Frontend Developer", CYAN),
    ("N3", "Name 3", "UI / Data Analytics", PURPLE),
]
cw3, cg3 = (CONTENT_W - 2 * 0.2) / 3, 0.2
for i, (ini, name, role, a) in enumerate(team):
    x = MARGIN + i * (cw3 + cg3)
    panel(s, x, 2.15, cw3, 3.55)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.03, 2.15, cw3 - 0.06, 0.08,
          fill=a, line=None, radius=0.5)
    shape(s, MSO_SHAPE.OVAL, x + cw3 / 2 - 0.55, 2.6, 1.1, 1.1, fill=None, line=a, line_w=2)
    tx(s, x + cw3 / 2 - 0.55, 2.6, 1.1, 1.1, ini, 30, a, bold=True,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx(s, x + 0.2, 3.95, cw3 - 0.4, 0.4, name, 20, TEXT, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x + 0.2, 4.45, cw3 - 0.4, 0.35, role, 12.5, MUTED, align=PP_ALIGN.CENTER)

tx(s, MARGIN, 6.15, CONTENT_W, 0.35,
   "\u2014  replace these placeholders with your team members  \u2014",
   12, MUTED, align=PP_ALIGN.CENTER)
footer(s, 3)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 · INTRODUCTION
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
content_header(s, "INTRODUCTION", "What is SentinalIQ?")

tx(s, MARGIN, 1.92, CONTENT_W, 1.05,
   "SentinalIQ is an Enterprise SIEM / SOC platform that gives security teams a "
   "single, real-time view of their environment — live dashboards, searchable "
   "logs, incident lifecycle management, endpoint & web-application monitoring, "
   "and threat intelligence, all secured with JWT authentication.",
   15, MUTED, wrap=True, line_spacing=1.25)

panels = [
    ("WHY WAS THIS PROJECT CHOSEN?", CYAN, [
        "Cybersecurity is one of the fastest-growing fields — real-world relevance",
        "Covers the full stack: React frontend, FastAPI backend & real-time systems",
        "Demonstrates core SOC concepts: SIEM, incident response, EDR, WAF, threat intel",
        "Hands-on practice with JWT auth, WebSockets & data visualization",
    ]),
    ("BACKGROUND OF THE PROBLEM", PURPLE, [
        "Organizations generate millions of security events every single day",
        "Analysts need one unified view to detect and respond to threats fast",
        "Traditional tools are fragmented, expensive and slow to react",
        "A modern, web-based, real-time SOC dashboard is the answer",
    ]),
]
pw = (CONTENT_W - 0.24) / 2
for i, (title, a, items) in enumerate(panels):
    x = MARGIN + i * (pw + 0.24)
    panel(s, x, 3.2, pw, 3.35)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.03, 3.2, pw - 0.06, 0.08,
          fill=a, line=None, radius=0.5)
    tx(s, x + 0.3, 3.42, pw - 0.6, 0.3, title, 12.5, a, bold=True, spacing=150)
    bullets(s, x + 0.3, 3.85, pw - 0.6, items, size=13, gap=10, accent=a)

footer(s, 4)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 · PROBLEM STATEMENT
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
content_header(s, "PROBLEM STATEMENT", "The Challenge", accent=RED)

problems = [
    ("Security data is scattered", " across endpoints, networks and cloud tools — no single source of truth"),
    ("Manual log analysis is slow", " and error-prone — it simply does not scale"),
    ("Slow detection & response", " leaves organizations exposed to breaches"),
    ("Alert fatigue", " — disconnected point tools generate too many noisy alerts"),
    ("Commercial SIEM platforms", " are expensive and complex for smaller teams"),
    ("No automated, real-time correlation", " of threats, endpoints and web traffic"),
]
box = s.shapes.add_textbox(Inches(MARGIN), Inches(2.05), Inches(CONTENT_W), Inches(1))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
for i, (head, rest) in enumerate(problems):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(14)
    p.line_spacing = 1.1
    g = p.add_run()
    g.text = "\u25B8  "
    g.font.size = Pt(16)
    g.font.bold = True
    g.font.color.rgb = RED
    g.font.name = FONT
    r1 = p.add_run()
    r1.text = head
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = TEXT
    r1.font.name = FONT
    r2 = p.add_run()
    r2.text = rest
    r2.font.size = Pt(16)
    r2.font.color.rgb = MUTED
    r2.font.name = FONT

footer(s, 5)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 · OBJECTIVES
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
content_header(s, "OBJECTIVES", "What We Set Out to Achieve")

objectives = [
    ("01", "Real-Time Visibility", "Build a live security dashboard with real-time metrics, charts and an activity feed.", CYAN),
    ("02", "Centralized Logs", "Aggregate security events from many sources into one searchable, filterable view.", PURPLE),
    ("03", "Automated Incident Lifecycle", "Manage incidents end-to-end: open → investigating → resolved.", PINK),
    ("04", "Faster Response", "Reduce detection-to-response time with instant alerts and notifications.", GREEN),
    ("05", "Secure Access", "Role-based JWT authentication for admins, analysts and SOC leads.", AMBER),
    ("06", "Actionable Intelligence", "Track IOCs & threat actors and surface analytics via dashboards.", CYAN),
]
ow = (CONTENT_W - 2 * 0.24) / 3
oh = 2.05
for i, (num, head, desc, a) in enumerate(objectives):
    x = MARGIN + (i % 3) * (ow + 0.24)
    y = 2.0 + (i // 3) * (oh + 0.28)
    panel(s, x, y, ow, oh)
    tx(s, x + 0.28, y + 0.2, 0.9, 0.5, num, 24, a, bold=True)
    tx(s, x + 0.28, y + 0.78, ow - 0.56, 0.4, head, 15.5, TEXT, bold=True)
    tx(s, x + 0.28, y + 1.2, ow - 0.56, 0.75, desc, 11.5, MUTED, wrap=True, line_spacing=1.15)

footer(s, 6)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 · EXISTING SYSTEM
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
content_header(s, "EXISTING SYSTEM", "Before SentinalIQ")

cols = [
    ("EXISTING SYSTEM", CYAN, [
        "Disconnected point tools & manual record-keeping",
        "Paper / email-based incident reporting",
        "No online or remote accessibility",
        "Reactive response — threats found after the fact",
        "No central database for security events",
    ]),
    ("LIMITATIONS", RED, [
        "Manual processes are time-consuming",
        "Data is not managed efficiently",
        "Lack of online accessibility",
        "High chances of human error",
        "No real-time visibility or alerting",
    ]),
]
cw7 = (CONTENT_W - 0.3) / 2
for i, (title, a, items) in enumerate(cols):
    x = MARGIN + i * (cw7 + 0.3)
    panel(s, x, 2.1, cw7, 4.4)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.03, 2.1, cw7 - 0.06, 0.09,
          fill=a, line=None, radius=0.5)
    tx(s, x + 0.3, 2.35, cw7 - 0.6, 0.35, title, 16, a, bold=True, spacing=200)
    bullets(s, x + 0.3, 2.95, cw7 - 0.6, items, size=14.5, gap=13, accent=a)

footer(s, 7)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 · TECHNOLOGIES USED
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
content_header(s, "TECHNOLOGIES USED", "The Tech Stack")

tech_rows = [
    ("Frontend", "React 19 · TypeScript · Vite · Recharts · Framer Motion", CYAN),
    ("Backend", "Python · FastAPI (REST + WebSocket)", PURPLE),
    ("Database", "SQLite · SQLAlchemy ORM", PINK),
    ("Authentication", "JWT · python-jose · passlib (bcrypt)", GREEN),
    ("Real-Time", "Native WebSockets · uvicorn[standard]", AMBER),
    ("UI / Styling", "Lucide React icons · CSS design-token system (dark & light)", CYAN),
    ("Server", "Uvicorn (ASGI)", PURPLE),
    ("IDE / Tools", "VS Code · Git / GitHub", PINK),
]
header_h, row_h = 0.5, 0.55
table_y = 1.95
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, table_y, CONTENT_W,
      header_h + len(tech_rows) * row_h, fill=BG2, line=EDGE, radius=0.04)
col_a, col_b = 3.3, CONTENT_W - 3.3
# header row
tx(s, MARGIN + 0.3, table_y + 0.09, col_a - 0.5, 0.32, "COMPONENT", 12.5, CYAN, bold=True, spacing=150)
tx(s, MARGIN + col_a + 0.3, table_y + 0.09, col_b - 0.5, 0.32, "TECHNOLOGY", 12.5, CYAN, bold=True, spacing=150)
shape(s, MSO_SHAPE.RECTANGLE, MARGIN + 0.02, table_y + header_h - 0.01, CONTENT_W - 0.04, 0.02, fill=EDGE, line=None)
for i, (comp, tech, a) in enumerate(tech_rows):
    y = table_y + header_h + i * row_h
    if i % 2 == 1:
        shape(s, MSO_SHAPE.RECTANGLE, MARGIN + 0.02, y, CONTENT_W - 0.04, row_h,
              fill=RGBColor(0x0D, 0x14, 0x26), line=None)
    shape(s, MSO_SHAPE.OVAL, MARGIN + 0.3, y + row_h / 2 - 0.05, 0.11, 0.11, fill=a, line=None)
    tx(s, MARGIN + 0.55, y, col_a - 0.7, row_h, comp, 13.5, TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    tx(s, MARGIN + col_a + 0.3, y, col_b - 0.45, row_h, tech, 13, MUTED, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 8)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 · FLOW OF PROJECT
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
content_header(s, "FLOW OF PROJECT", "Architecture & Data Flow")

flow_box(s, 0.7, 2.75, 2.0, 1.15, "User / Browser", "React SPA · HTTPS")
shape(s, MSO_SHAPE.RIGHT_ARROW, 2.74, 3.18, 0.4, 0.28, fill=CYAN, line=None)
flow_box(s, 3.16, 2.75, 2.45, 1.15, "React Frontend", "Vite · TypeScript · Recharts", CYAN)
shape(s, MSO_SHAPE.RIGHT_ARROW, 5.65, 3.18, 0.4, 0.28, fill=CYAN, line=None)
flow_box(s, 6.07, 2.75, 2.75, 1.15, "FastAPI Backend", "REST API · WebSocket · JWT Auth", PURPLE)
shape(s, MSO_SHAPE.RIGHT_ARROW, 8.86, 3.18, 0.4, 0.28, fill=CYAN, line=None)
flow_box(s, 9.28, 2.75, 2.3, 1.15, "SQLite Database", "SQLAlchemy ORM · Seeded data", GREEN)

# WebSocket feedback loop
line(s, 7.44, 3.9, 7.44, 5.3, color=CYAN, dash=True, head=True)
line(s, 7.44, 5.3, 4.38, 5.3, color=CYAN, dash=True)
line(s, 4.38, 5.3, 4.38, 3.9, color=CYAN, dash=True, head=True)
panel(s, 4.6, 4.98, 3.9, 0.64, fill=BG2, line_c=CYAN)
tx(s, 4.6, 4.98, 3.9, 0.64, "WebSocket · Real-time alerts & live updates",
   11.5, CYAN, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

tx(s, MARGIN, 6.1, CONTENT_W, 0.6,
   "The browser calls the FastAPI backend over REST; the backend persists to SQLite and "
   "streams real-time alerts & events back to the UI over WebSockets.",
   12.5, MUTED, align=PP_ALIGN.CENTER, wrap=True, line_spacing=1.2)
footer(s, 9)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 · MAIN LOGIC
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
content_header(s, "MAIN LOGIC", "Core Implementation")

code_a = (
    "# auth.py — JWT token creation\n"
    "def create_access_token(data: dict,\n"
    "                        expires_delta=None) -> str:\n"
    "    to_encode = data.copy()\n"
    "    expire = datetime.now(timezone.utc) + (\n"
    "        expires_delta or\n"
    "        timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))\n"
    "    to_encode.update({\"exp\": expire,\n"
    "                      \"jti\": generate_uuid()})\n"
    "    return jwt.encode(to_encode, SECRET_KEY,\n"
    "                      algorithm=ALGORITHM)\n"
)

code_b = (
    "# server.py — dashboard metrics\n"
    "@app.get(\"/api/dashboard/stats\")\n"
    "def dashboard_stats(\n"
    "        payload=Depends(get_current_user),\n"
    "        db: Session = Depends(get_db)):\n"
    "    user_id = payload[\"sub\"]\n"
    "    total = db.query(Incident).filter(\n"
    "        Incident.user_id == user_id).count()\n"
    "    critical = db.query(Incident).filter(\n"
    "        Incident.severity == \"critical\").count()\n"
    "    return {\"total_incidents\": total,\n"
    "            \"critical\": critical}\n"
)

code_c = (
    "# server.py — real-time WebSocket\n"
    "@app.websocket(\"/ws\")\n"
    "async def websocket_endpoint(ws: WebSocket):\n"
    "    await ws.accept()\n"
    "    await ws.send_json({\n"
    "        \"type\": \"connected\",\n"
    "        \"message\": \"SentinalIQ live stream\"})\n"
    "    while True:\n"
    "        data = await ws.receive_text()\n"
    "        if json.loads(data).get(\"type\") == \"ping\":\n"
    "            await ws.send_json({\"type\": \"pong\"})\n"
)

cw10 = (CONTENT_W - 2 * 0.16) / 3
code_panel(s, MARGIN, 1.95, cw10, 4.6, "auth.py · JWT signing", code_a, PINK)
code_panel(s, MARGIN + (cw10 + 0.16), 1.95, cw10, 4.6, "server.py · Dashboard API", code_b, CYAN)
code_panel(s, MARGIN + 2 * (cw10 + 0.16), 1.95, cw10, 4.6, "server.py · Live WebSocket", code_c, GREEN)

tx(s, MARGIN, 6.75, CONTENT_W, 0.3,
   "Key implementation snippets from the SentinalIQ backend.", 11, MUTED, align=PP_ALIGN.CENTER)
footer(s, 10)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 · FUTURE SCOPE
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
content_header(s, "FUTURE SCOPE", "Where SentinalIQ Goes Next")

future = [
    ("AI / ML anomaly detection", "Model-driven threat scoring and predictive alerts beyond rule-based engines."),
    ("Cloud deployment", "Containerized Docker + Kubernetes rollout on AWS / Azure / GCP."),
    ("Mobile application", "On-the-go SOC monitoring with push notifications."),
    ("Third-party integrations", "Connect external EDR, firewall, cloud and ticketing systems."),
    ("Multi-tenant & multi-language", "SaaS-style tenant isolation and localized UI support."),
    ("SOAR automation", "Automated playbooks that respond to incidents without human action."),
]
fw = (CONTENT_W - 0.24) / 2
for i, (head, desc) in enumerate(future):
    x = MARGIN + (i % 2) * (fw + 0.24)
    y = 2.0 + (i // 2) * 1.55
    panel(s, x, y, fw, 1.38)
    tx(s, x + 0.28, y + 0.16, fw - 0.56, 0.35, head, 15, TEXT, bold=True)
    tx(s, x + 0.28, y + 0.56, fw - 0.56, 0.7, desc, 11.5, MUTED, wrap=True, line_spacing=1.15)

footer(s, 11)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 · THANK YOU
# ═════════════════════════════════════════════════════════════════════════════
s = slide()
bg(s)
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 0.09, SLIDE_H, fill=CYAN, line=None)
shape(s, MSO_SHAPE.RECTANGLE, 0.09, 0, 0.09, SLIDE_H, fill=PURPLE, line=None)
shape(s, MSO_SHAPE.RECTANGLE, 0.18, 0, 0.09, SLIDE_H, fill=PINK, line=None)

tx(s, 1.0, 2.15, 11.33, 1.1, "Thank You", 64, TEXT, bold=True, align=PP_ALIGN.CENTER)
tx(s, 1.0, 3.5, 11.33, 0.4, "Questions are welcome", 18, MUTED, align=PP_ALIGN.CENTER)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 5.57, 4.25, 2.2, 0.05, fill=CYAN, line=None, radius=0.5)

tx(s, 1.0, 4.7, 11.33, 0.4, "SentinalIQ — Enterprise SIEM / SOC Platform", 16, CYAN, bold=True, align=PP_ALIGN.CENTER)
tx(s, 1.0, 5.2, 11.33, 0.3, "School of Computing  ·  Sponsored by ASIA Charitable Trust", 12, MUTED, align=PP_ALIGN.CENTER)
tx(s, 1.0, 5.6, 11.33, 0.3, "Tech Cast · Project Presentation", 11, MUTED, align=PP_ALIGN.CENTER, spacing=250)

# ── Save ─────────────────────────────────────────────────────────────────────
OUT = "SentinalIQ_Presentation.pptx"
prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides._sldIdLst)} slides")
